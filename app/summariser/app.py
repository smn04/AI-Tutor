from flask import Flask, render_template, request, redirect, url_for, Blueprint
from PyPDF2 import PdfReader
from pptx import Presentation
from summa.summarizer import summarize

summariser_bp = Blueprint('summariser', __name__, template_folder='templates')


@summariser_bp.route('/')
def summariser_index():
    return render_template('summariser_index.html')

@summariser_bp.route('/file_summariser', methods=['GET', 'POST'])
def file_summariser():
    if request.method == 'POST':
        file_type = request.form.get('file_type')
        if file_type == 'pdf':
            return redirect(url_for('summariser.pdf_summariser'))
        elif file_type == 'ppt':
            return redirect(url_for('summariser.ppt_summariser'))
        else:
            return redirect(url_for('summariser.text_summariser'))
    return render_template('summariser_index.html')

@summariser_bp.route('/pdf_summariser', methods=['GET', 'POST'])
def pdf_summariser():
    if request.method == 'POST':
        file = request.files['file']
        summary_ratio = float(request.form['summary_ratio']) / 100.0

        with file.stream as pdf_file:
            text = ""
            pdf_reader = PdfReader(pdf_file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text()
            summary = summarize(text, ratio=summary_ratio)
            summary_text=summary
            return render_template('summary_result.html', summary_text=summary_text)

    return render_template('pdf_summariser.html')

@summariser_bp.route('/ppt_summariser', methods=['GET', 'POST'])
def ppt_summariser():
    if request.method == 'POST':
        file = request.files['file']
        summary_ratio = float(request.form['summary_ratio']) / 100.0

        with file.stream as pptx_file:
            presentation = Presentation(pptx_file)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text

            summary = summarize(text, ratio=summary_ratio)
            summary_text=summary

            return render_template('summary_result.html', summary_text=summary_text)

    return render_template('ppt_summariser.html')

@summariser_bp.route('/text_summariser', methods=['GET', 'POST'])
def text_summariser():
    if request.method == 'POST':
        user_text = request.form['user_text']
        summary_ratio = float(request.form['summary_ratio']) / 100.0

        summary = summarize(user_text, ratio=summary_ratio)
        summary_text=summary
        return render_template('summary_result.html', summary_text=summary_text)

    return render_template('text_summariser.html')


